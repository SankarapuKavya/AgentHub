import os, json, uuid, re, mimetypes, time, zipfile, shutil, io, threading
from repo_parser import is_repo_output, parse_repo, write_repo, build_tree, REPO_SYSTEM_ADDENDUM
from datetime import datetime
from pathlib import Path
from flask import Flask, render_template, request, jsonify, send_file
import requests

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024

BASE_DIR   = Path(__file__).parent
DATA_FILE  = BASE_DIR / "data" / "store.json"
OUTPUT_DIR = Path(r"C:\Users\ganesh.r.vennapusa\OneDrive - Accenture\Ganesh\AM Outputs")
UPLOAD_DIR = BASE_DIR / "uploads"
RUNS_DIR   = BASE_DIR / "runs"
TMP_DIR    = BASE_DIR / "tmp"

for d in [DATA_FILE.parent, UPLOAD_DIR, RUNS_DIR, TMP_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# ── Active pipeline run sessions ──────────────────────────────────────────
# run_id -> {
#   "status": "running" | "hil_waiting" | "done" | "error",
#   "steps_done": [...],          # completed step records
#   "hil_event": threading.Event, # set when user responds
#   "hil_response": {},           # user's approve/revise decision
#   "hil_info": {},               # what we're pausing on
#   "final_output": str,
#   "output_files": [],
#   "error": str,
#   "run_rec": {}                 # full run record when done
# }
SESSIONS = {}
SESSIONS_LOCK = threading.Lock()

# ── Store ─────────────────────────────────────────────────────────────────
def load():
    if DATA_FILE.exists():
        return json.loads(DATA_FILE.read_text())
    return {"agents": [], "super_agents": [], "files": [], "runs": [], "settings": {}}

def save(data):
    DATA_FILE.write_text(json.dumps(data, indent=2))

# ── LLM call ──────────────────────────────────────────────────────────────
_last_usage = [{}]  # thread-local token tracking

def call_llm(system_prompt, user_parts, settings, retries=3):
    api_key  = settings.get("api_key", "")
    api_url  = settings.get("api_url", "").rstrip("/")
    model    = settings.get("model", "")
    provider = settings.get("provider", "openai").lower()

    if not api_key or not api_url or not model:
        raise ValueError("LLM not configured. Open Settings in the sidebar.")

    # Inject repo output instructions when prompt is repo-type
    REPO_RE = r"\b(repo(sitory)?|project folder|code project|codebase|generate.*project|scaffold|boilerplate|folder structure)\b"
    if re.search(REPO_RE, system_prompt, re.I):
        system_prompt = system_prompt + "\n\n" + REPO_SYSTEM_ADDENDUM

    # Hard limits to stay well under Groq's payload size limit
    # System prompt: keep full. Each file: cap at 6000 chars. Total user content: cap at 24000 chars.
    MAX_FILE_CHARS  = 8_000   # per file
    MAX_TOTAL_CHARS = 40_000  # total user content

    parts_text = []
    for p in user_parts:
        if p["type"] == "text" and p.get("text", "").strip():
            text = p["text"]
            # Cap previous-agent output at 12000 chars too
            if text.startswith("Previous") and len(text) > 30_000:
                # Only truncate very large outputs — repo code must not be cut
                text = text[:30_000] + "\n\n[... output truncated ...]"
            parts_text.append(text)
        elif p["type"] == "file":
            file_content = p["content"]
            if len(file_content) > MAX_FILE_CHARS:
                file_content = file_content[:MAX_FILE_CHARS] + f"\n\n[... file truncated at {MAX_FILE_CHARS} chars to fit context ...]"
            parts_text.append(f"\n--- File: {p['name']} ---\n{file_content}\n--- End ---\n")

    user_content = "\n\n".join(parts_text).strip() or "Process this."

    # Final safety cap on total payload
    if len(user_content) > MAX_TOTAL_CHARS:
        user_content = user_content[:MAX_TOTAL_CHARS] + "\n\n[... content truncated ...]"

    # Cap system prompt but ensure it stays complete (don't cut mid-sentence)
    if len(system_prompt) > 8_000:
        system_prompt = system_prompt[:8_000] + "\n[system prompt truncated]"

    headers = {"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}

    if provider == "gemini":
        url = f"{api_url}/models/{model}:generateContent?key={api_key}"
        payload = {"contents": [{"parts": [{"text": f"{system_prompt}\n\n{user_content}"}]}]}
        for attempt in range(retries):
            try:
                r = requests.post(url, json=payload, timeout=180); r.raise_for_status()
                gd = r.json()
                meta = gd.get("usageMetadata",{})
                _last_usage[0] = {"prompt_tokens": meta.get("promptTokenCount",0), "completion_tokens": meta.get("candidatesTokenCount",0), "total_tokens": meta.get("totalTokenCount",0)}
                return gd["candidates"][0]["content"]["parts"][0]["text"]
            except requests.exceptions.HTTPError:
                if r.status_code == 429 and attempt < retries - 1:
                    wait = (attempt + 1) * 20
                    print(f"Groq 429 rate limit hit — waiting {wait}s before retry {attempt+2}/{retries}")
                    time.sleep(wait)
                    continue
                if r.status_code == 413:
                    raise ValueError(
                        "Payload too large (413). Your input content is too long. "
                        "Try: shorter text input, smaller files, or split the task across multiple agents."
                    )
                if r.status_code == 429:
                    raise ValueError(
                        "Rate limit hit (429). You've exceeded the free tier limits. "
                        "Wait 1 minute and try again, or switch to a different model like "
                        "llama-3.1-8b-instant or meta-llama/llama-4-scout-17b-16e-instruct."
                    )
                raise
    else:
        endpoint = api_url if api_url.endswith("/chat/completions") else f"{api_url}/chat/completions"
        payload = {
            "model": model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content}
            ],
            "max_tokens": 8192, "temperature": 0.7
        }
        for attempt in range(retries):
            try:
                r = requests.post(endpoint, headers=headers, json=payload, timeout=180); r.raise_for_status()
                rd = r.json()
                # Store last token usage in thread-local-ish via return tuple
                usage = rd.get("usage", {})
                _last_usage[0] = {"prompt_tokens": usage.get("prompt_tokens",0), "completion_tokens": usage.get("completion_tokens",0), "total_tokens": usage.get("total_tokens",0)}
                return rd["choices"][0]["message"]["content"]
            except requests.exceptions.HTTPError:
                if r.status_code == 429 and attempt < retries - 1:
                    wait = (attempt + 1) * 20
                    print(f"Groq 429 rate limit hit — waiting {wait}s before retry {attempt+2}/{retries}")
                    time.sleep(wait)
                    continue
                if r.status_code == 413:
                    raise ValueError(
                        "Payload too large (413). Your input content is too long. "
                        "Try: shorter text input, smaller files, or split the task across multiple agents."
                    )
                if r.status_code == 429:
                    raise ValueError(
                        "Rate limit hit (429). You've exceeded the free tier limits. "
                        "Wait 1 minute and try again, or switch to a different model like "
                        "llama-3.1-8b-instant or meta-llama/llama-4-scout-17b-16e-instruct."
                    )
                raise

# ── Format detection ──────────────────────────────────────────────────────
FORMAT_RE = {
    "xlsx": r"\b(excel|xlsx|spreadsheet)\b",
    "docx": r"\b(word|docx|word\s*doc(?:ument)?)\b",
    "pptx": r"\b(powerpoint|pptx|presentation|slides?)\b",
    "pdf":  r"\b(pdf)\b",
    "csv":  r"\b(csv|comma[\s\-]separated)\b",
    "json": r"\b(json)\b",
    "html": r"\b(html|web\s*page)\b",
    "md":   r"\b(markdown|\.md)\b",
    "py":   r"\b(python\s*(?:script|code)|\.py)\b",
    "js":   r"\b(javascript\s*(?:script|file)|\.js)\b",
}

def detect_formats(prompt):
    p_lower = prompt.lower()
    # Repo detection — check for explicit repo/project/folder keywords
    REPO_RE = r"\b(repo(sitory)?|project folder|code project|code base|codebase|generate.*project|scaffold|boilerplate|folder structure|directory structure|file structure)\b"
    if re.search(REPO_RE, p_lower, re.I):
        return ["repo"]
    f = [k for k, p in FORMAT_RE.items() if re.search(p, prompt, re.I)]
    return f if f else ["txt"]

def write_file(content, fmt, filepath):
    # "repo" format is handled separately — this shouldn't be called with fmt=="repo"
    # but guard just in case
    if fmt == "repo":
        filepath.write_text(content, encoding="utf-8")
        return filepath
    TEXT = {"txt", "md", "html", "csv", "json", "py", "js"}
    if fmt in TEXT:
        filepath.write_text(content, encoding="utf-8")
    elif fmt == "docx":
        try:
            from docx import Document
            doc = Document()
            for line in content.split("\n"):
                s = line.strip()
                if s.startswith("### "): doc.add_heading(s[4:], level=3)
                elif s.startswith("## "): doc.add_heading(s[3:], level=2)
                elif s.startswith("# "): doc.add_heading(s[2:], level=1)
                elif s.startswith(("- ", "* ")): doc.add_paragraph(s[2:], style="List Bullet")
                elif s: doc.add_paragraph(s)
            doc.save(filepath)
        except ImportError:
            filepath = filepath.with_suffix(".txt"); filepath.write_text(content, encoding="utf-8")
    elif fmt == "xlsx":
        try:
            import openpyxl
            wb = openpyxl.Workbook(); ws = wb.active
            for i, line in enumerate(content.split("\n"), 1):
                if not line.strip(): continue
                cols = [c.strip().strip('"') for c in re.split(r'\t|,(?=(?:[^"]*"[^"]*")*[^"]*$)', line)]
                for j, val in enumerate(cols, 1): ws.cell(row=i, column=j, value=val)
            wb.save(filepath)
        except ImportError:
            filepath = filepath.with_suffix(".txt"); filepath.write_text(content, encoding="utf-8")
    elif fmt == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation()
            for block in re.split(r'\n(?=#{1,2} |Slide \d)', content):
                lines = [l for l in block.strip().split("\n") if l.strip()]
                if not lines: continue
                sl = prs.slides.add_slide(prs.slide_layouts[1])
                sl.shapes.title.text = lines[0].lstrip("#").strip()
                sl.placeholders[1].text = "\n".join(lines[1:])
            prs.save(filepath)
        except ImportError:
            filepath = filepath.with_suffix(".txt"); filepath.write_text(content, encoding="utf-8")
    elif fmt == "pdf":
        try:
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import inch
            styles = getSampleStyleSheet(); story = []
            for line in content.split("\n"):
                if line.strip():
                    story.append(Paragraph(line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), styles["Normal"]))
                else:
                    story.append(Spacer(1, 0.1 * inch))
            SimpleDocTemplate(str(filepath)).build(story)
        except ImportError:
            filepath = filepath.with_suffix(".txt"); filepath.write_text(content, encoding="utf-8")
    else:
        filepath.write_text(content, encoding="utf-8")
    return filepath

def read_file_content(path, mime):
    suffix = path.suffix.lower()
    try:
        if "text" in mime or suffix in (".txt", ".md", ".csv", ".json", ".py", ".js", ".html", ".xml"):
            return path.read_text(encoding="utf-8", errors="ignore")
        elif suffix == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(path) as pdf:
                    return "\n".join(p.extract_text() or "" for p in pdf.pages)
            except:
                try:
                    import pypdf
                    reader = pypdf.PdfReader(str(path))
                    return "\n".join(p.extract_text() or "" for p in reader.pages)
                except: return f"[PDF: {path.name} — extraction unavailable]"
        elif suffix == ".docx":
            try:
                from docx import Document
                doc = Document(path); return "\n".join(p.text for p in doc.paragraphs)
            except: return f"[DOCX: {path.name} — extraction unavailable]"
        elif suffix in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True); lines = []
                for ws in wb.worksheets:
                    lines.append(f"Sheet: {ws.title}")
                    for row in ws.iter_rows(values_only=True):
                        lines.append(",".join(str(c) if c is not None else "" for c in row))
                return "\n".join(lines)
            except: return f"[Excel: {path.name} — extraction unavailable]"
        else:
            return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        return f"[Error reading {path.name}: {e}]"

# ══════════════════════════════════════════════════════════════════════════
# Settings
# ══════════════════════════════════════════════════════════════════════════
@app.route("/api/settings", methods=["GET"])
def api_get_settings():
    d = load(); s = d.get("settings", {}); masked = dict(s)
    if masked.get("api_key"):
        k = masked["api_key"]
        masked["api_key_masked"] = k[:6] + "•" * max(0, len(k) - 10) + k[-4:] if len(k) > 10 else "••••••"
    return jsonify(masked)

@app.route("/api/settings", methods=["POST"])
def api_save_settings():
    d = load(); b = request.json
    s = {"provider": b.get("provider", "openai"), "api_url": b.get("api_url", "").strip(), "model": b.get("model", "").strip()}
    if b.get("api_key"): s["api_key"] = b["api_key"].strip()
    elif d.get("settings", {}).get("api_key"): s["api_key"] = d["settings"]["api_key"]
    d["settings"] = s; save(d); return jsonify({"ok": True})

# ══════════════════════════════════════════════════════════════════════════
# Files
# ══════════════════════════════════════════════════════════════════════════
@app.route("/api/files", methods=["GET"])
def api_files(): return jsonify(load()["files"])

@app.route("/api/files/upload", methods=["POST"])
def api_upload():
    if "file" not in request.files: return jsonify({"error": "No file"}), 400
    f = request.files["file"]; fid = str(uuid.uuid4())
    mime = f.content_type or mimetypes.guess_type(f.filename)[0] or "application/octet-stream"
    dest = UPLOAD_DIR / fid; dest.mkdir(); fp = dest / f.filename; f.save(fp)
    d = load()
    rec = {"id": fid, "name": f.filename, "mime": mime, "size": fp.stat().st_size, "uploaded_at": datetime.now().isoformat()}
    d["files"].append(rec); save(d); return jsonify(rec), 201

@app.route("/api/files/<fid>", methods=["DELETE"])
def api_del_file(fid):
    d = load(); d["files"] = [f for f in d["files"] if f["id"] != fid]; save(d)
    p = UPLOAD_DIR / fid
    if p.exists(): shutil.rmtree(p)
    return jsonify({"ok": True})

# ══════════════════════════════════════════════════════════════════════════
# Agents
# ══════════════════════════════════════════════════════════════════════════
@app.route("/api/agents", methods=["GET"])
def api_agents(): return jsonify(load()["agents"])

@app.route("/api/agents", methods=["POST"])
def api_create_agent():
    d = load(); b = request.json
    a = {"id": str(uuid.uuid4()), "name": b["name"], "description": b.get("description", ""),
         "prompt": b.get("prompt", ""), "type": b.get("type", "llm"), "created_at": datetime.now().isoformat()}
    d["agents"].append(a); save(d); return jsonify(a), 201

@app.route("/api/agents/<aid>", methods=["PUT"])
def api_update_agent(aid):
    d = load(); b = request.json
    for a in d["agents"]:
        if a["id"] == aid:
            for k in ("name", "description", "prompt", "type"):
                if k in b: a[k] = b[k]
            save(d); return jsonify(a)
    return jsonify({"error": "Not found"}), 404

@app.route("/api/agents/<aid>", methods=["DELETE"])
def api_del_agent(aid):
    d = load(); d["agents"] = [a for a in d["agents"] if a["id"] != aid]; save(d)
    return jsonify({"ok": True})

# ══════════════════════════════════════════════════════════════════════════
# Pipelines
# ══════════════════════════════════════════════════════════════════════════
@app.route("/api/super_agents", methods=["GET"])
def api_supers(): return jsonify(load()["super_agents"])

@app.route("/api/super_agents", methods=["POST"])
def api_create_super():
    d = load(); b = request.json
    sa = {"id": str(uuid.uuid4()), "name": b["name"], "description": b.get("description", ""),
          "flow": b.get("flow", []), "created_at": datetime.now().isoformat()}
    d["super_agents"].append(sa); save(d); return jsonify(sa), 201

@app.route("/api/super_agents/<sid>", methods=["PUT"])
def api_update_super(sid):
    d = load(); b = request.json
    for sa in d["super_agents"]:
        if sa["id"] == sid:
            for k in ("name", "description", "flow"):
                if k in b: sa[k] = b[k]
            save(d); return jsonify(sa)
    return jsonify({"error": "Not found"}), 404

@app.route("/api/super_agents/<sid>", methods=["DELETE"])
def api_del_super(sid):
    d = load(); d["super_agents"] = [s for s in d["super_agents"] if s["id"] != sid]; save(d)
    return jsonify({"ok": True})

# ══════════════════════════════════════════════════════════════════════════
# PIPELINE STEPS — HIL is a step type inside the flow list
# flow item: {"type": "agent", "agent_id": "..."} 
#         or {"type": "hil", "id": "hil-xxx", "label": "Review checkpoint"}
# ══════════════════════════════════════════════════════════════════════════

# ══════════════════════════════════════════════════════════════════════════
# RUN — step-by-step execution with true HIL blocking
# ══════════════════════════════════════════════════════════════════════════

def _run_pipeline_thread(run_id, sid, user_text, file_parts, settings):
    """Runs in a background thread. Blocks at HIL steps waiting for user input."""
    d = load()
    sa = next((s for s in d["super_agents"] if s["id"] == sid), None)
    agent_map = {a["id"]: a for a in d["agents"]}
    flow = sa.get("flow", [])  # list of step dicts

    run_dir = RUNS_DIR / run_id; run_dir.mkdir(exist_ok=True)
    ts_str = datetime.now().strftime("%Y%m%d_%H%M%S")
    sa_slug = re.sub(r'[^\w\-]', '_', sa["name"])

    steps_done = []
    all_paths = []
    current_output = user_text

    sess = SESSIONS[run_id]

    try:
        for idx, step in enumerate(flow):
            step_num = idx + 1

            # ── HIL STEP ─────────────────────────────────────────────────
            if step.get("type") == "hil":
                # ── HIL LOOP: pause → user reviews → approve or revise N-1 → re-pause
                hil_revision_count = 0
                while True:
                    # Find all completed LLM steps (for revision target dropdown)
                    llm_steps_done = [s for s in steps_done if s.get("step_type") == "llm"]

                    # Find the immediately preceding LLM step (default re-run target)
                    prev_llm_step = llm_steps_done[-1] if llm_steps_done else None

                    hil_info = {
                        "step": step_num,
                        "label": step.get("label", "Review checkpoint"),
                        "current_output": current_output,
                        "revision_count": hil_revision_count,
                        "prev_step_name": prev_llm_step["agent_name"] if prev_llm_step else None,
                        "llm_steps": [{"step": s["step"], "agent_name": s["agent_name"], "agent_id": s.get("agent_id")} for s in llm_steps_done]
                    }

                    # Signal pause and wait
                    with SESSIONS_LOCK:
                        sess["status"] = "hil_waiting"
                        sess["hil_info"] = hil_info
                        sess["hil_event"].clear()
                        sess["hil_response"] = None

                    got = sess["hil_event"].wait(timeout=7200)
                    if not got:
                        with SESSIONS_LOCK:
                            sess["status"] = "error"; sess["error"] = "HIL timed out"
                        return

                    response = sess["hil_response"]

                    # ── APPROVED: exit loop, continue pipeline ────────────
                    if response["action"] == "approve":
                        steps_done.append({
                            "step": step_num, "step_type": "hil",
                            "agent_name": step.get("label", "Review checkpoint"),
                            "output": "[Approved by user]", "files": [], "hil_approved": True,
                            "revision_count": hil_revision_count
                        })
                        with SESSIONS_LOCK:
                            sess["status"] = "running"
                            sess["steps_done"] = list(steps_done)
                        break  # exit HIL loop, move to next pipeline step

                    # ── REVISE: re-run the previous LLM step, then loop ──
                    elif response["action"] == "revise":
                        feedback = response.get("feedback", "")
                        target_agent_id = response.get("agent_id", "")

                        # Resolve which agent to re-run
                        target_agent = agent_map.get(target_agent_id) if target_agent_id else None
                        prev_step_num = None
                        if not target_agent and prev_llm_step:
                            target_agent = agent_map.get(prev_llm_step["agent_id"])
                            prev_step_num = prev_llm_step["step"]
                        elif target_agent:
                            # find the step number for this agent
                            matching = next((s for s in llm_steps_done if s.get("agent_id")==target_agent_id), None)
                            prev_step_num = matching["step"] if matching else None

                        with SESSIONS_LOCK:
                            sess["status"] = "running"
                            sess["current_step_running"] = prev_step_num
                            sess["steps_done"] = list(steps_done)

                        if target_agent:
                            hil_revision_count += 1
                            revision_parts = [
                                {"type": "text", "text": (
                                    f"Your previous output:\n{current_output}\n\n"
                                    f"The reviewer has requested changes:\n{feedback}\n\n"
                                    f"Please revise your output accordingly. "
                                    f"This is revision #{hil_revision_count}."
                                )}
                            ]
                            # Include original file_parts on first revision so agent has full context
                            if idx == 0 or hil_revision_count == 1:
                                pass  # file_parts already in context via current_output
                            try:
                                revised_output = call_llm(target_agent["prompt"], revision_parts, settings)
                            except Exception as e:
                                revised_output = f"[Revision error: {e}]"

                            # Save revised files
                            fmts = detect_formats(target_agent["prompt"])
                            agent_slug = re.sub(r"[^\w\-]", "_", target_agent["name"])
                            rev_tag = f"_rev{hil_revision_count}"
                            revised_files = []
                            if "repo" in fmts or is_repo_output(revised_output):
                                try:
                                    proj_name = f"step{prev_step_num}_{agent_slug}{rev_tag}"
                                    parsed_files = parse_repo(revised_output)
                                    if parsed_files:
                                        zip_path, _ = write_repo(parsed_files, run_dir, proj_name)
                                        tree = build_tree(parsed_files)
                                        revised_files.append({"filename": f"{proj_name}.zip","path": str(zip_path),"format": "repo","repo_files": [f["path"] for f in parsed_files],"tree": tree,"file_count": len(parsed_files)})
                                        all_paths.append(zip_path)
                                except Exception as e:
                                    revised_files.append({"filename": "repo_error.txt","path": f"Error:{e}","format":"txt"})
                            else:
                                for fmt in fmts:
                                    fname = f"step{prev_step_num}_{agent_slug}{rev_tag}.{fmt}"
                                    fpath = run_dir / fname
                                    try:
                                        write_file(revised_output, fmt, fpath)
                                        revised_files.append({"filename": fname, "path": str(fpath), "format": fmt})
                                        all_paths.append(fpath)
                                    except Exception as e:
                                        revised_files.append({"filename": fname, "path": f"Error:{e}", "format": fmt})

                            # Update current_output so the next HIL review sees the revision
                            current_output = revised_output

                            # Add a REVISION record so the frontend can show it
                            revision_record = {
                                "step": step_num, "step_type": "hil_revision",
                                "agent_name": f"Revision {hil_revision_count} — {target_agent['name']}",
                                "agent_id": target_agent["id"],
                                "output": revised_output,
                                "files": revised_files,
                                "feedback": feedback,
                                "hil_approved": False,
                                "revision_num": hil_revision_count
                            }
                            steps_done.append(revision_record)

                            with SESSIONS_LOCK:
                                sess["current_step_running"] = None
                                sess["steps_done"] = list(steps_done)
                            # Loop back to re-pause at HIL for user to review again

                        else:
                            # No agent found — just continue
                            steps_done.append({
                                "step": step_num, "step_type": "hil",
                                "agent_name": step.get("label", "Review checkpoint"),
                                "output": "[No prior LLM agent found to revise]", "files": []
                            })
                            with SESSIONS_LOCK:
                                sess["status"] = "running"
                                sess["steps_done"] = list(steps_done)
                            break

                continue  # move to next flow step

            # ── NORMAL AGENT STEP ─────────────────────────────────────────
            if step.get("type") != "agent":
                continue

            agent = agent_map.get(step.get("agent_id"))
            if not agent:
                steps_done.append({"step": step_num, "step_type": "llm", "agent_id": None,
                                    "agent_name": "Unknown", "output": "[Agent not found]", "files": []})
                with SESSIONS_LOCK:
                    sess["steps_done"] = list(steps_done)
                continue

            # Build user parts
            step_parts = []
            if idx == 0:
                if user_text: step_parts.append({"type": "text", "text": f"User Input:\n{user_text}"})
                step_parts.extend(file_parts)
            else:
                step_parts.append({"type": "text", "text": f"Previous step output:\n{current_output}"})

            # Signal "running this step"
            with SESSIONS_LOCK:
                sess["current_step_running"] = step_num

            # For repo steps, add explicit instruction to output COMPLETE files
            _prompt = agent["prompt"]
            _fmts_check = detect_formats(_prompt)
            if "repo" in _fmts_check or is_repo_output(""):
                pass  # REPO_SYSTEM_ADDENDUM already injected in call_llm

            try:
                llm_out = call_llm(_prompt, step_parts, settings)
                step_tokens = dict(_last_usage[0])
            except Exception as e:
                llm_out = f"[ERROR: {e}]"
                step_tokens = {}

            fmts = detect_formats(_prompt)
            agent_slug = re.sub(r'[^\w\-]', '_', agent["name"])
            step_files = []

            if "repo" in fmts or is_repo_output(llm_out):
                # ── REPO OUTPUT: parse into folder structure ──────────────
                try:
                    proj_name = f"step{step_num}_{agent_slug}"
                    parsed_files = parse_repo(llm_out)
                    if parsed_files:
                        proj_dir, written_paths = write_repo(parsed_files, run_dir, proj_name)
                        tree = build_tree(parsed_files)
                        # Record the project folder for the step (for UI display)
                        step_files.append({
                            "filename": proj_name,
                            "path": str(proj_dir),
                            "format": "repo",
                            "repo_files": [f["path"] for f in parsed_files],
                            "tree": tree,
                            "file_count": len(parsed_files)
                        })
                        # Add each individual file to all_paths — outer zip bundles them
                        for wp in written_paths:
                            all_paths.append(wp)
                    else:
                        # LLM said repo but output wasn't parseable — save as txt
                        fpath = run_dir / f"step{step_num}_{agent_slug}.txt"
                        fpath.write_text(llm_out, encoding="utf-8")
                        step_files.append({"filename": fpath.name, "path": str(fpath), "format": "txt"})
                        all_paths.append(fpath)
                except Exception as e:
                    step_files.append({"filename": "repo_error.txt", "path": f"Error: {e}", "format": "txt"})
            else:
                # ── NORMAL FILE OUTPUT ────────────────────────────────────
                for fmt in fmts:
                    fname = f"step{step_num}_{agent_slug}.{fmt}"
                    fpath = run_dir / fname
                    try:
                        write_file(llm_out, fmt, fpath)
                        step_files.append({"filename": fname, "path": str(fpath), "format": fmt})
                        all_paths.append(fpath)
                    except Exception as e:
                        step_files.append({"filename": fname, "path": f"Error: {e}", "format": fmt})

            steps_done.append({
                "step": step_num, "step_type": "llm",
                "agent_id": agent["id"], "agent_name": agent["name"],
                "output": llm_out, "files": step_files,
                "tokens": step_tokens
            })
            current_output = llm_out
            total_tok = sum(s.get("tokens",{}).get("total_tokens",0) for s in steps_done)
            with SESSIONS_LOCK:
                sess["steps_done"] = list(steps_done)
                sess["current_step_running"] = None
                sess["total_tokens"] = total_tok

        # ── FINALIZE ──────────────────────────────────────────────────────
        output_files = []
        try:
            OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
            if len(all_paths) > 0:
                zname = f"{sa_slug}_{ts_str}.zip"
                zpath = OUTPUT_DIR / zname
                with zipfile.ZipFile(zpath, "w", zipfile.ZIP_DEFLATED) as outer_zf:
                    for p in all_paths:
                        if not p.exists():
                            continue
                        # Preserve relative path inside run_dir so folder structure is kept
                        try:
                            arcname = p.relative_to(run_dir)
                        except ValueError:
                            arcname = p.name
                        outer_zf.write(p, arcname)
                output_files = [{"filename": zname, "path": str(zpath), "format": "zip"}]
        except Exception as e:
            output_files = [{"filename": "error", "path": str(e), "format": "txt"}]

        run_rec = {
            "id": run_id, "pipeline_id": sid, "pipeline_name": sa["name"],
            "started_at": sess.get("started_at"), "finished_at": datetime.now().isoformat(),
            "status": "completed", "input_text": user_text,
            "input_files": [p["name"] for p in file_parts],
            "steps": steps_done, "output_files": output_files,
            "final_output": current_output
        }
        dd = load(); dd.setdefault("runs", []).append(run_rec); save(dd)

        with SESSIONS_LOCK:
            sess["status"] = "done"
            sess["final_output"] = current_output
            sess["output_files"] = output_files
            sess["steps_done"] = steps_done
            sess["run_rec"] = run_rec

    except Exception as e:
        with SESSIONS_LOCK:
            sess["status"] = "error"
            sess["error"] = str(e)


@app.route("/api/pipelines/<sid>/run", methods=["POST"])
def api_run(sid):
    d = load()
    sa = next((s for s in d["super_agents"] if s["id"] == sid), None)
    if not sa: return jsonify({"error": "Pipeline not found"}), 404
    if not sa.get("flow"): return jsonify({"error": "Pipeline has no steps"}), 400

    settings = d.get("settings", {})
    user_text = request.form.get("input_text", "").strip()
    sel_fids = request.form.getlist("file_ids")

    file_parts = []
    for fid in sel_fids:
        rec = next((f for f in d["files"] if f["id"] == fid), None)
        if not rec: continue
        fp = UPLOAD_DIR / fid / rec["name"]
        if fp.exists():
            file_parts.append({"type": "file", "name": rec["name"],
                               "content": read_file_content(fp, rec["mime"]), "mime": rec["mime"]})

    if "upload" in request.files:
        uf = request.files["upload"]; tmp = TMP_DIR / uf.filename; uf.save(tmp)
        mime = uf.content_type or mimetypes.guess_type(uf.filename)[0] or "text/plain"
        file_parts.append({"type": "file", "name": uf.filename,
                           "content": read_file_content(tmp, mime), "mime": mime})

    if not user_text and not file_parts:
        return jsonify({"error": "Provide text or files"}), 400

    run_id = str(uuid.uuid4())
    session = {
        "status": "running",
        "started_at": datetime.now().isoformat(),
        "pipeline_name": sa["name"],
        "steps_done": [],
        "current_step_running": None,
        "hil_event": threading.Event(),
        "hil_response": None,
        "hil_info": None,
        "final_output": None,
        "output_files": [],
        "error": None,
        "run_rec": None
    }
    with SESSIONS_LOCK:
        SESSIONS[run_id] = session

    t = threading.Thread(target=_run_pipeline_thread,
                         args=(run_id, sid, user_text, file_parts, settings), daemon=True)
    t.start()
    return jsonify({"run_id": run_id})


@app.route("/api/runs/<rid>/status", methods=["GET"])
def api_run_status(rid):
    with SESSIONS_LOCK:
        sess = SESSIONS.get(rid)
    if not sess:
        # Check persisted runs
        d = load()
        r = next((r for r in d.get("runs", []) if r["id"] == rid), None)
        if r: return jsonify({"status": "done", "run": r})
        return jsonify({"error": "Not found"}), 404

    out = {
        "status": sess["status"],
        "steps_done": sess["steps_done"],
        "current_step_running": sess.get("current_step_running"),
        "final_output": sess.get("final_output"),
        "output_files": sess.get("output_files", []),
        "error": sess.get("error"),
        "total_tokens": sess.get("total_tokens", 0),
    }
    if sess["status"] == "hil_waiting":
        out["hil_info"] = sess.get("hil_info")
    if sess["status"] == "done" and sess.get("run_rec"):
        out["run"] = sess["run_rec"]
    return jsonify(out)


@app.route("/api/runs/<rid>/hil", methods=["POST"])
def api_hil_respond(rid):
    """User approves or sends revision feedback."""
    with SESSIONS_LOCK:
        sess = SESSIONS.get(rid)
    if not sess: return jsonify({"error": "Run not found"}), 404
    if sess["status"] != "hil_waiting":
        return jsonify({"error": "Run is not waiting for HIL input"}), 400

    b = request.json
    action = b.get("action")  # "approve" or "revise"
    if action not in ("approve", "revise"):
        return jsonify({"error": "action must be 'approve' or 'revise'"}), 400

    with SESSIONS_LOCK:
        sess["hil_response"] = {
            "action": action,
            "feedback": b.get("feedback", ""),
            "agent_id": b.get("agent_id"),   # which agent to re-run
        }
        sess["hil_event"].set()

    return jsonify({"ok": True})


@app.route("/api/runs/<rid>/steps/<int:step>/repo")
def dl_step_repo(rid, step):
    """Download the repo zip for a specific step."""
    with SESSIONS_LOCK:
        sess = SESSIONS.get(rid)
    steps = sess["steps_done"] if sess else []
    if not steps:
        d = load(); r = next((r for r in d.get("runs", []) if r["id"] == rid), None)
        steps = r["steps"] if r else []
    s = next((s for s in steps if s["step"] == step), None)
    if not s: return jsonify({"error": "Step not found"}), 404
    f = next((f for f in s.get("files", []) if f.get("format") == "repo"), None)
    if not f: return jsonify({"error": "No repo output for this step"}), 404
    p = Path(f["path"])
    if not p.exists(): return jsonify({"error": "File not found on disk"}), 404
    proj_name = p.stem  # filename without .zip
    return send_file(str(p), as_attachment=True, download_name=f"{proj_name}.zip", mimetype="application/zip")


@app.route("/api/runs/<rid>/steps/<int:step>/download/<fmt>")
def dl_step(rid, step, fmt):
    with SESSIONS_LOCK:
        sess = SESSIONS.get(rid)
    steps = sess["steps_done"] if sess else []
    if not steps:
        d = load(); r = next((r for r in d.get("runs", []) if r["id"] == rid), None)
        steps = r["steps"] if r else []
    s = next((s for s in steps if s["step"] == step), None)
    if not s: return jsonify({"error": "Step not found"}), 404
    f = next((f for f in s.get("files", []) if f["format"] == fmt), None)
    if not f: return jsonify({"error": "File not found"}), 404
    p = Path(f["path"])
    if not p.exists():
        tmp = TMP_DIR / f["filename"]; write_file(s.get("output", ""), fmt, tmp); p = tmp
    return send_file(str(p), as_attachment=True, download_name=f["filename"])


@app.route("/api/runs/<rid>/download")
def dl_run(rid):
    with SESSIONS_LOCK:
        sess = SESSIONS.get(rid)
    steps = sess["steps_done"] if sess else []
    if not steps:
        d = load(); r = next((r for r in d.get("runs", []) if r["id"] == rid), None)
        steps = r["steps"] if r else []
    run_dir_path = RUNS_DIR / rid
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for s in steps:
            for f in s.get("files", []):
                p = Path(f["path"])
                if not p.exists():
                    continue
                if p.is_dir():
                    # Repo folder — add all files inside preserving structure
                    for child in p.rglob("*"):
                        if child.is_file():
                            try:
                                arcname = child.relative_to(run_dir_path)
                            except ValueError:
                                arcname = child.relative_to(p.parent)
                            zf.write(child, arcname)
                else:
                    try:
                        arcname = p.relative_to(run_dir_path)
                    except ValueError:
                        arcname = p.name
                    zf.write(p, arcname)
    buf.seek(0)
    return send_file(buf, as_attachment=True, download_name=f"outputs_{rid[:8]}.zip", mimetype="application/zip")


@app.route("/api/runs", methods=["GET"])
def api_runs():
    d = load(); q = request.args.get("q", "").lower(); runs = d.get("runs", [])
    if q:
        runs = [r for r in runs if q in r.get("pipeline_name", "").lower()
                or any(q in s.get("agent_name", "").lower() for s in r.get("steps", []))]
    return jsonify(sorted(runs, key=lambda r: r.get("started_at", ""), reverse=True))


@app.route("/api/runs/<rid>", methods=["GET"])
def api_run_detail(rid):
    d = load(); r = next((r for r in d.get("runs", []) if r["id"] == rid), None)
    return jsonify(r) if r else (jsonify({"error": "Not found"}), 404)


@app.route("/")
def index(): return render_template("index.html")


if __name__ == "__main__":
    app.run(debug=True, port=5000, threaded=True)
